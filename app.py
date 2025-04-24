import os
os.environ["TORCH_USE_RTLD_GLOBAL"] = "YES"
GOOGLE_API_KEY = "AIzaSyBV44ZygJe0bcNus8tIuUCtoYpjyT0LdZw"  # add your GOOGLE API key here
os.environ["GOOGLE_API_KEY"] = GOOGLE_API_KEY
import json
import torch
import numpy as np
from PIL import Image
import pandas as pd
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
import logging
import tempfile
import io
import shutil
import streamlit as st
import requests
from bs4 import BeautifulSoup
import re
import time

# Core LlamaIndex imports
from llama_index.core import (
    VectorStoreIndex, 
    SimpleDirectoryReader,
    PromptTemplate,
    ServiceContext,
    Settings
)
from llama_index.core.node_parser import SentenceSplitter
from llama_index.llms.huggingface import HuggingFaceLLM
from llama_index.llms.gemini import Gemini
from llama_index.embeddings.huggingface import HuggingFaceEmbedding
from llama_index.embeddings.gemini import GeminiEmbedding
from llama_index.core.schema import Document
from llama_index.core.query_engine import RetrieverQueryEngine

# Document loaders
from llama_index.readers.file import PDFReader, DocxReader

# For OCR processing
import pytesseract
from pdf2image import convert_from_path
import fitz  # PyMuPDF
import cv2

# For presentations
from pptx import Presentation

# For tables
import tabula
import camelot

# For DOCX processing
import docx

# Set up logging
logging.basicConfig(
    filename='app.log',
    level=logging.INFO, 
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

class DocumentProcessor:
    """Enhanced document processor to handle multiple file formats with special focus on PDF and DOCX processing."""
    
    def __init__(
        self,
        embedding_model: str = "BAAI/bge-small-en-v1.5",  # Free, open-source
        llm_model: str = "models/gemini-1.5-flash", #"StabilityAI/stablelm-tuned-alpha-3b",  # Free, open-source
        chunk_size: int = 1024,
        chunk_overlap: int = 100,  # New hyperparameter
        use_ocr: bool = True,
        similarity_top_k: int = 5,  # New hyperparameter
        similarity_cutoff: float = 0.7,  # New hyperparameter
        crawl_links: bool = True,  # New hyperparameter
        temperature: float = 0.2,  # New hyperparameter for LLM
        max_links_to_crawl: int = 5,  # New hyperparameter
        structured_data_threshold: int = 5  # New hyperparameter for data detection
    ):
        #self.data_dir = data_dir
        self.use_ocr = use_ocr
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.similarity_top_k = similarity_top_k
        self.similarity_cutoff = similarity_cutoff
        self.crawl_links = crawl_links
        self.temperature = temperature
        self.max_links_to_crawl = max_links_to_crawl
        self.structured_data_threshold = structured_data_threshold
        
        # Create data directory if it doesn't exist
        #os.makedirs(data_dir, exist_ok=True)
        
        # Initialize embedding model
        logger.info("Initializing embedding model...")
        Settings.embed_model=GeminiEmbedding(
        model_name=embedding_model, api_key=GOOGLE_API_KEY
        )
        #Settings.embed_model = HuggingFaceEmbedding(model_name=embedding_model)
        
        # Initialize LLM
        logger.info("Initializing LLM...")
        self._initialize_llm(llm_model)
        
        # Create service context
        Settings.chunk_size = chunk_size
        Settings.llm = self.llm
        
        # Initialize node parser for text chunking
        self.node_parser = SentenceSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
        
        # Storage for all extracted links
        self.extracted_links = []
        
        # Storage for loaded dataframes (CSV/Excel)
        self.loaded_dataframes = {}
    
    def _initialize_llm(self, model_name: str):
        """Initialize the LLM with appropriate settings."""
        # Define the system prompt for the LLM
        system_prompt = """<|SYSTEM|># StableLM Tuned (Alpha version)
        - StableLM is a helpful and harmless open-source AI language model developed by StabilityAI.
        - StableLM is excited to be able to help the user, but will refuse to do anything that could be considered harmful to the user.
        - StableLM is more than just an information source, StableLM is also able to write poetry, short stories, and make jokes.
        - StableLM will refuse to participate in anything that could harm a human.
        
        Your task is to answer questions based solely on the provided document context. If the information cannot be found in the documents, acknowledge that you don't know rather than making up information. Always be concise and precise in your answers, citing the source when possible.
        """

        # Create the prompt template
        query_wrapper_prompt = PromptTemplate("<|USER|>{query_str}<|ASSISTANT|>")
        
        self.llm = Gemini(
            model=model_name,  # Assuming `model_name` is one of GEMINI_MODELS
            temperature=self.temperature,
            max_tokens=512,
            generation_config=None,  # Optional, pass if you have a custom config
            safety_settings=None,    # Optional, pass if needed
            callback_manager=None,   # Optional, pass if using callbacks
            api_key=None,            # Optional, include your API key if required
            api_base=None,           # Optional, specify if you're using a custom endpoint
            transport=None,          # Optional, default should work in most cases
            model_name=model_name,   # Redundant if same as `model`, but included just in case
            default_headers=None,    # Optional
            request_options=None,    # Optional
            do_sample=True,
            pad_token_id=0
        )
    
    def _extract_text_from_pdf_page(self, page) -> str:
        """Extract regular text from a PDF page using PyMuPDF."""
        return page.get_text()
    
    def _extract_tables_from_pdf_page(self, page, page_num: int, pdf_path: str) -> List[str]:
        """Extract tables from a PDF page using both PyMuPDF's table finder and tabula/camelot."""
        table_texts = []
        
        # 1. Try PyMuPDF's table detection first
        # try:
        #     tables = page.find_tables()
        #     if tables and tables.tables:
        #         for i, table in enumerate(tables.tables):
        #             rows = []
        #             for row in table.rows:
        #                 row_data = [cell.text for cell in row.cells]
        #                 rows.append(" | ".join(row_data))
        #             table_text = "\n".join(rows)
        #             table_texts.append(f"Table {i+1} on page {page_num+1}:\n{table_text}")
        # except Exception as e:
        #     logger.warning(f"PyMuPDF table extraction failed on page {page_num+1}: {e}")
        
        # 2. Try tabula for more advanced table extraction (Java-based, works well for many tables)
        # try:
        #     tables = tabula.read_pdf(pdf_path, pages=page_num+1)
        #     for i, table in enumerate(tables):
        #         if not table.empty:
        #             table_text = f"Table {i+1} detected by tabula on page {page_num+1}:\n{table.to_string()}"
        #             table_texts.append(table_text)
        # except Exception as e:
        #     logger.warning(f"Tabula table extraction failed on page {page_num+1}: {e}")
        
        # 3. Try camelot as a final option (better for complex tables but slower)
        # if not table_texts:
        try:
            tables = camelot.read_pdf(pdf_path, pages=str(page_num+1))
            if len(tables) > 0:
                for i, table in enumerate(tables):
                    table_df = table.df
                    if not table_df.empty:
                        table_text = f"Table {i+1} detected by camelot on page {page_num+1}:\n{table_df.to_string()}"
                        table_texts.append(table_text)
        except Exception as e:
            logger.warning(f"Camelot table extraction failed on page {page_num+1}: {e}")
        
        return table_texts

    def _extract_links_from_pdf_page(self, page) -> List[Dict[str, Any]]:
        """Extract hyperlinks from a PDF page and return as structured data."""
        links_data = []
        try:
            for link in page.get_links():
                uri = link.get("uri")
                if uri:
                    if uri not in [l["link"] for l in self.extracted_links]:
                        link_data = {
                            "type": "hyperlink",
                            "link": uri,
                        }
                        links_data.append(link_data)
                        self.extracted_links.append(link_data)
        except Exception as e:
            logger.warning(f"Failed to extract links from PDF page: {e}")
        
        return links_data

    def _crawl_links(self, urls: List[str]) -> List[Dict[str, Any]]:
        """Crawl a list of URLs and return their content."""
        crawled_content = []
        links_crawled=0
        for url in urls:
            try:
                # Skip non-HTTP URLs
                if not url.startswith(('http://', 'https://')):
                    continue
                    
                logger.info(f"Crawling URL: {url}")

                if links_crawled >= self.max_links_to_crawl:
                    logger.info(f"Reached maximum links to crawl ({self.max_links_to_crawl})")
                    break                
                
                # Fetch the page with timeout
                response = requests.get(url, timeout=10)
                response.raise_for_status()
                
                # Parse and clean the content
                soup = BeautifulSoup(response.text, 'html.parser')
                
                # Remove scripts and styles
                for element in soup(['script', 'style', 'nav', 'footer', 'header']):
                    element.decompose()
                    
                # Get clean text
                text = soup.get_text()
                lines = (line.strip() for line in text.splitlines())
                chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
                clean_text = '\n'.join(chunk for chunk in chunks if chunk)
                
                if clean_text:
                    crawled_content.append({
                        "type": "link_content",
                        "link": url,
                        "content": clean_text
                    })
                    
                # Be polite with delay
                links_crawled += 1
                time.sleep(1)
                
            except Exception as e:
                logger.warning(f"Failed to crawl URL {url}: {e}")
        
        return crawled_content
    
    # def _crawl_links(self) -> List[Dict[str, Any]]:
    #     """Crawl extracted links to gather additional content."""
    #     logger.info(f"Crawling {len(self.extracted_links)} extracted links...")
    #     all_contents = []

    #     if not self.crawl_links or not self.extracted_links:
    #         return all_contents

    #     links_crawled = 0
    #     for link_data in self.extracted_links:
    #         url = link_data["link"]

    #         if not url.startswith(('http://', 'https://')):
    #             continue

    #         if links_crawled >= self.max_links_to_crawl:
    #             logger.info(f"Reached maximum links to crawl ({self.max_links_to_crawl})")
    #             break

    #         try:
    #             logger.info(f"Crawling link: {url}")
    #             response = requests.get(url, timeout=10)
    #             response.raise_for_status()

    #             soup = BeautifulSoup(response.text, 'html.parser')
    #             for script in soup(["script", "style"]):
    #                 script.extract()

    #             text = soup.get_text()
    #             lines = (line.strip() for line in text.splitlines())
    #             chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
    #             cleaned_text = '\n'.join(chunk for chunk in chunks if chunk)

    #             if cleaned_text:
    #                 all_contents.append({
    #                     "type": "link_content",
    #                     "link": url,
    #                     "content": cleaned_text
    #                 })

    #                 links_crawled += 1
    #                 time.sleep(1)

    #         except Exception as e:
    #             logger.warning(f"Failed to crawl link {url}: {e}")

    #     logger.info(f"Successfully crawled {links_crawled} links")
    #     return all_contents


    def _get_images_from_pdf_page(self, page) -> List[Image.Image]:
        """Extract images from a PDF page using PyMuPDF."""
        images = []
        try:
            for img_index, img in enumerate(page.get_images(full=True)):
                xref = img[0]
                base_image = page.parent.extract_image(xref)
                image_bytes = base_image["image"]
                
                # Load the image with PIL
                try:
                    image = Image.open(io.BytesIO(image_bytes))
                    images.append(image)
                except Exception as e:
                    logger.warning(f"Failed to open image {img_index} from PDF: {e}")
        except Exception as e:
            logger.warning(f"Failed to extract images from PDF page: {e}")
        
        return images
    
    def _ocr_image(self, image: Image.Image) -> str:
        """Process an image with OCR to extract text."""
        try:
            # Improve image for OCR
            # Convert to grayscale
            gray_image = cv2.cvtColor(np.array(image), cv2.COLOR_RGB2GRAY)
            
            # Apply threshold to get black and white image
            _, binary_image = cv2.threshold(gray_image, 150, 255, cv2.THRESH_BINARY)
            
            # Convert back to PIL image for Tesseract
            binary_pil = Image.fromarray(binary_image)
            
            # Apply OCR with improved settings
            text = pytesseract.image_to_string(
                binary_pil, 
                config='--psm 1 --oem 3'  # Page segmentation mode 1: Auto page segmentation with OSD
            )
            return text.strip()
        except Exception as e:
            logger.error(f"OCR processing error: {e}")
            return ""
    
    def _process_pdf_with_ocr(self, file_path: str) -> List[Dict[str, Any]]:
        """Process PDF with OCR to extract text, tables, links, and web content."""
        logger.info(f"Processing PDF with OCR: {file_path}")
        all_contents = []

        # Open the PDF
        doc = fitz.open(file_path)

        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            page_contents = []

            # 1. Extract regular text
            try:
                text = self._extract_text_from_pdf_page(page)
                if text.strip():  
                    page_contents.append({
                        "type": "text",
                        "content": text,
                        "page": page_num + 1
                    })
                    logger.debug(f"Page {page_num+1}: Extracted {len(text)} characters of text")
            except Exception as e:
                logger.warning(f"Failed to extract text from page {page_num+1}: {e}")
            
            # 2. Extract tables from the page
            try:
                table_texts =self. _extract_tables_from_pdf_page(page, page_num, file_path)
                for table_text in table_texts:
                    page_contents.append({
                        "type": "table",
                        "content": table_text,
                        "page": page_num + 1
                    })
                    logger.debug(f"Page {page_num+1}: Extracted table with {len(table_text)} characters")
            except Exception as e:
                logger.warning(f"Failed to extract tables from page {page_num+1}: {e}")
            # 3. Extract links
            try:
                links = self._extract_links_from_pdf_page(page)
                if links:
                    links_text = f"Links on page {page_num+1}:\n" + "\n".join([l["link"] for l in links if "link" in l])
                    page_contents.append({
                        "type": "hyperlinks",
                        "content": links_text,
                        "page": page_num + 1
                    })
                    logger.debug(f"Page {page_num+1}: Extracted {len(links)} links")
                    
                    if self.crawl_links:
                        link_list=[i.get('link') for i in links]
                        crawled_content = self._crawl_links(link_list)
                        all_contents.extend(crawled_content)

            except Exception as e:
                logger.warning(f"Failed to extract links from page {page_num+1}: {e}")
            

            # 4. OCR on full page if text is minimal or OCR explicitly enabled
            if self.use_ocr and (len(text.strip()) < 100 or not text.strip()):
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
                    pix.save(tmp.name)
                    tmp_path = tmp.name

                try:
                    img = Image.open(tmp_path)
                    ocr_text = self._ocr_image(img)
                    if ocr_text.strip():
                        page_contents.append({
                            "type": "ocr_text",
                            "content": f"OCR text from page {page_num+1}:\n{ocr_text}",
                            "page": page_num + 1
                        })
                finally:
                    os.unlink(tmp_path)

            # 5. OCR on embedded images
            if self.use_ocr:
                for img_index, img in enumerate(page.get_images(full=True)):
                    try:
                        xref = img[0]
                        base_image = doc.extract_image(xref)
                        image_bytes = base_image["image"]

                        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
                            with open(tmp.name, 'wb') as f:
                                f.write(image_bytes)

                            img = Image.open(tmp.name)
                            ocr_text = self._ocr_image(img)

                            if ocr_text.strip():
                                page_contents.append({
                                    "type": "embedded_image_ocr",
                                    "content": f"Embedded image {img_index+1} OCR from page {page_num+1}:\n{ocr_text}",
                                    "page": page_num + 1
                                })

                            os.unlink(tmp.name)
                    except Exception as e:
                        logger.warning(f"Failed to process embedded image {img_index} on page {page_num+1}: {e}")

            all_contents.extend(page_contents)

        return all_contents

    def _extract_hyperlinks_from_docx(self,doc) -> List[str]:
        """Extract all hyperlink URLs from a DOCX document."""
        hyperlinks = []
        
        # Method 1: Extract hyperlinks from document relationships
        if hasattr(doc, 'part') and hasattr(doc.part, 'rels'):
            for rel in doc.part.rels.values():
                if "hyperlink" in rel.reltype:
                    target = rel.target_ref
                    if target.startswith(('http://', 'https://')):
                        hyperlinks.append(target)
        
        # Method 2: Alternative way to find hyperlinks in paragraphs
        for paragraph in doc.paragraphs:
            for run in paragraph.runs:
                # Check if run has hyperlink by examining XML
                if run._element.xpath('.//w:hyperlink'):
                    hyperlink_elements = run._element.xpath('.//w:hyperlink')
                    for hyperlink in hyperlink_elements:
                        if 'id' in hyperlink.attrib:
                            rel_id = hyperlink.attrib['{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id']
                            if rel_id in doc.part.rels:
                                target = doc.part.rels[rel_id].target_ref
                                if target.startswith(('http://', 'https://')):
                                    hyperlinks.append(target)
        
        # Deduplicate links
        return list(set(hyperlinks))    

    def _process_docx(self, file_path: str) -> List[Dict[str, Any]]:
        """Process DOCX file to extract text, tables, images, and hyperlinks."""
        logger.info(f"Processing DOCX file: {file_path}")
        all_contents = []
        
        try:
            doc = docx.Document(file_path)
            
            # 1. Extract all text
            full_text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
            if full_text:
                all_contents.append({
                    "type": "text",
                    "content": full_text
                })
            
            # 2. Extract hyperlinks
            hyperlinks = self._extract_hyperlinks_from_docx(doc)
            if hyperlinks:
                # Add hyperlinks to output
                all_contents.append({
                    "type": "hyperlinks",
                    "content": "Document links:\n" + "\n".join(hyperlinks)
                })
                
                # Crawl links if enabled
                if self.crawl_links:
                    crawled_content = self._crawl_links(hyperlinks)
                    all_contents.extend(crawled_content)
            
            # 3. Extract tables
            for i, table in enumerate(doc.tables):
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(" | ".join(row_data))
                
                if table_data:
                    all_contents.append({
                        "type": "table",
                        "content": f"Table {i+1}:\n" + "\n".join(table_data)
                    })
            
            # 4. Extract images (OCR if enabled)
            if self.use_ocr:
                if hasattr(doc, 'part'):
                    temp_dir = tempfile.mkdtemp()
                    try:
                        image_count = 0
                        for rel in doc.part.rels.values():
                            if "image" in rel.target_ref.lower():
                                try:
                                    image_path = os.path.join(temp_dir, f"image_{image_count}.png")
                                    with open(image_path, "wb") as f:
                                        f.write(rel.target_part.blob)
                                    
                                    ocr_text = _ocr_image(Image.open(image_path))
                                    if ocr_text.strip():
                                        all_contents.append({
                                            "type": "image_ocr",
                                            "content": f"Image {image_count+1} text:\n{ocr_text}"
                                        })
                                    image_count += 1
                                except Exception as e:
                                    logger.warning(f"Failed to process image {image_count}: {e}")
                    finally:
                        shutil.rmtree(temp_dir)
                    
        except Exception as e:
            logger.error(f"Error processing DOCX {file_path}: {e}", exc_info=True)
        
        return all_contents
    
    # def _process_docx(self, file_path: str) -> List[Dict[str, Any]]:
    #     """Process DOCX file with full extraction of text, tables, images, and links."""
    #     logger.info(f"Processing DOCX file: {file_path}")
    #     all_contents = []

    #     try:
    #         doc = docx.Document(file_path)
    #         paragraphs_text = []
    #         hyperlinks = []

    #         # Extract text and hyperlinks
    #         for para in doc.paragraphs:
    #             paragraphs_text.append(para.text)

    #             for run in para.runs:
    #                 if hasattr(run, '_element') and run._element.xpath('.//w:hyperlink'):
    #                     for hyperlink in run._element.xpath('.//w:hyperlink'):
    #                         if 'r:id' in hyperlink.attrib:
    #                             rid = hyperlink.attrib['{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id']
    #                             rels = doc.part.rels
    #                             if rid in rels:
    #                                 target = rels[rid].target_ref
    #                                 link_text = "".join([t.text for t in hyperlink.xpath('.//w:t')])
    #                                 link_info = f"Link: {link_text} -> {target}"
    #                                 hyperlinks.append(link_info)
    #                                 if link_info not in self.extracted_links:
    #                                     self.extracted_links.append(link_info)

    #         if paragraphs_text:
    #             all_contents.append({
    #                 "type": "text",
    #                 "content": "\n".join(paragraphs_text)
    #             })

    #         if hyperlinks:
    #             all_contents.append({
    #                 "type": "hyperlinks",
    #                 "content": f"Links in document:\n" + "\n".join(hyperlinks)
    #             })

    #             if self.crawl_links:
    #                 crawled_link_contents = self._crawl_links()
    #                 all_contents.extend(crawled_link_contents)

    #         # Extract tables
    #         for i, table in enumerate(doc.tables):
    #             rows = []
    #             for row in table.rows:
    #                 row_content = [cell.text for cell in row.cells]
    #                 rows.append(" | ".join(row_content))

    #             table_text = f"Table {i+1}:\n" + "\n".join(rows)
    #             all_contents.append({
    #                 "type": "table",
    #                 "content": table_text
    #             })

    #         # Extract and OCR images
    #         if self.use_ocr:
    #             temp_dir = tempfile.mkdtemp()
    #             try:
    #                 image_count = 0
    #                 for rel in doc.part.rels.values():
    #                     if "image" in rel.target_ref:
    #                         try:
    #                             image_blob = rel.target_part.blob
    #                             image_filename = os.path.join(temp_dir, f"image_{image_count}.png")
    #                             with open(image_filename, "wb") as f:
    #                                 f.write(image_blob)

    #                             img = Image.open(image_filename)
    #                             ocr_text = self._ocr_image(img)

    #                             if ocr_text.strip():
    #                                 all_contents.append({
    #                                     "type": "embedded_image_ocr",
    #                                     "content": f"Embedded image {image_count+1} OCR:\n{ocr_text}"
    #                                 })

    #                             image_count += 1
    #                         except Exception as e:
    #                             logger.warning(f"Failed to process embedded image {image_count} in DOCX: {e}")
    #             finally:
    #                 shutil.rmtree(temp_dir)

    #     except Exception as e:
    #         logger.error(f"Error processing DOCX {file_path}: {e}", exc_info=True)

    #     return all_contents

    
    def _process_image_with_ocr(self, file_path: str) -> str:
        """Process image with OCR to extract text."""
        logger.info(f"Processing image with OCR: {file_path}")
        try:
            # Load image
            image = Image.open(file_path)
            
            # Apply OCR
            return self._ocr_image(image)
        except Exception as e:
            logger.error(f"Error processing image {file_path}: {e}")
            return ""
    
    def _process_excel(self, file_path: str) -> List[str]:
        """Process Excel file."""
        logger.info(f"Processing Excel file: {file_path}")
        all_text = []
        
        # Read all sheets
        excel_file = pd.ExcelFile(file_path, engine='openpyxl')
        
        # Store the dataframe for potential querying
        filename = os.path.basename(file_path)
        self.loaded_dataframes[filename] = {}
        
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name, engine='openpyxl')
            self.loaded_dataframes[filename][sheet_name] = df
            all_text.append(f"Sheet: {sheet_name}\n{df.to_string()}")
        
        return all_text
    
    def _process_powerpoint(self, file_path: str) -> List[Dict[str, Any]]:
        """Process PowerPoint file to extract text, images (OCR), and hyperlinks."""
        logger.info(f"Processing PowerPoint file: {file_path}")
        all_contents = []
        
        try:
            prs = Presentation(file_path)
            
            for i, slide in enumerate(prs.slides):
                slide_text = f"Slide {i+1}:\n"
                hyperlinks = []

                # Extract text and hyperlinks from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text += shape.text + "\n"
                    
                    if hasattr(shape, "_element"):
                        link_elements = shape._element.xpath('.//a:hlinkClick')
                        for link in link_elements:
                            rid = link.attrib.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                            if rid and rid in slide.part.rels:
                                target = slide.part.rels[rid].target_ref
                                if target.startswith("http"):
                                    link_info = {'type': 'hyperlink', 'link': target}
                                    if link_info not in hyperlinks and link_info not in self.extracted_links:
                                        hyperlinks.append(link_info)
                                        self.extracted_links.append(link_info)

                # Add slide text
                if slide_text.strip() != f"Slide {i+1}:":
                    all_contents.append({
                        "type": "text",
                        "content": slide_text,
                        "slide": i+1
                    })

                # Add hyperlinks
                if hyperlinks:
                    links_text = f"Links on slide {i+1}:\n" + "\n".join([h["link"] for h in hyperlinks])
                    all_contents.append({
                        "type": "hyperlinks",
                        "content": links_text,
                        "slide": i+1
                    })

                    # Crawl the links if enabled
                    if self.crawl_links:
                        crawled = self._crawl_links(hyperlinks)
                        all_contents.extend(crawled)

                # OCR on images
                if self.use_ocr:
                    temp_dir = tempfile.mkdtemp()
                    try:
                        image_count = 0
                        for shape in slide.shapes:
                            if hasattr(shape, "image"):
                                try:
                                    image_bytes = shape.image.blob
                                    image_path = os.path.join(temp_dir, f"slide_{i+1}_img_{image_count}.png")
                                    with open(image_path, "wb") as f:
                                        f.write(image_bytes)
                                    
                                    ocr_text = self._ocr_image(Image.open(image_path))
                                    if ocr_text.strip():
                                        all_contents.append({
                                            "type": "embedded_image_ocr",
                                            "content": f"Image {image_count+1} OCR from slide {i+1}:\n{ocr_text}",
                                            "slide": i+1
                                        })
                                    image_count += 1
                                except Exception as e:
                                    logger.warning(f"Failed to process image in slide {i+1}: {e}")
                    finally:
                        shutil.rmtree(temp_dir)
        
        except Exception as e:
            logger.error(f"Error processing PowerPoint {file_path}: {e}", exc_info=True)
        
        return all_contents

    # def _process_powerpoint(self, file_path: str) -> List[Dict[str, Any]]:
    #     """Process PowerPoint file."""
    #     logger.info(f"Processing PowerPoint file: {file_path}")
    #     all_contents = []
        
    #     try:
    #         prs = Presentation(file_path)
            
    #         # Process each slide
    #         for i, slide in enumerate(prs.slides):
    #             slide_text = f"Slide {i+1}:\n"
    #             hyperlinks = []
                
    #             # Extract text content
    #             for shape in slide.shapes:
    #                 if hasattr(shape, "text") and shape.text:
    #                     slide_text += shape.text + "\n"
                    
    #                 # Extract hyperlinks if available
    #                 if hasattr(shape, "_element") and shape._element.xpath('.//a:hlinkClick'):
    #                     for link in shape._element.xpath('.//a:hlinkClick'):
    #                         if hasattr(link, 'attrib') and 'r:id' in link.attrib:
    #                             rid = link.attrib['{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id']
    #                             if hasattr(slide, 'rels') and rid in slide.rels:
    #                                 target = slide.rels[rid].target_ref
    #                                 link_text = shape.text if hasattr(shape, "text") else "[Shape]"
    #                                 link_info = f"Link: {link_text} -> {target}"
    #                                 hyperlinks.append(link_info)
    #                                 if link_info not in self.extracted_links:
    #                                     self.extracted_links.append(link_info)
                
    #             # Add text content
    #             if slide_text.strip() != f"Slide {i+1}:":
    #                 all_contents.append({
    #                     "type": "text",
    #                     "content": slide_text,
    #                     "slide": i+1
    #                 })
                
    #             # Add hyperlinks
    #             if hyperlinks:
    #                 all_contents.append({
    #                     "type": "hyperlinks",
    #                     "content": f"Links on slide {i+1}:\n" + "\n".join(hyperlinks),
    #                     "slide": i+1
    #                 })
                
    #             # Extract and process images if OCR is enabled
    #             if self.use_ocr:
    #                 temp_dir = tempfile.mkdtemp()
    #                 try:
    #                     image_count = 0
    #                     for shape in slide.shapes:
    #                         if hasattr(shape, "image"):
    #                             try:
    #                                 # Try to get image data
    #                                 image_bytes = shape.image.blob
    #                                 if image_bytes:
    #                                     image_filename = os.path.join(temp_dir, f"slide_{i+1}_image_{image_count}.png")
    #                                     with open(image_filename, "wb") as f:
    #                                         f.write(image_bytes)
                                        
    #                                     # Process with OCR
    #                                     img = Image.open(image_filename)
    #                                     ocr_text = self._ocr_image(img)
                                        
    #                                     if ocr_text.strip():
    #                                         all_contents.append({
    #                                             "type": "embedded_image_ocr",
    #                                             "content": f"Image {image_count+1} OCR from slide {i+1}:\n{ocr_text}",
    #                                             "slide": i+1
    #                                         })
                                        
    #                                     image_count += 1
    #                             except Exception as e:
    #                                 logger.warning(f"Failed to process image in slide {i+1}: {e}")
    #                 finally:
    #                     # Clean up
    #                     shutil.rmtree(temp_dir)
            
    #         if self.crawl_links:
    #             crawled_contents = self._crawl_links()
    #             all_contents.extend(crawled_contents)

    #     except Exception as e:
    #         logger.error(f"Error processing PowerPoint {file_path}: {e}")
        
    #     return all_contents
    
    def _process_json(self, file_path: str) -> str:
        """Process JSON file."""
        logger.info(f"Processing JSON file: {file_path}")
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            
            # Store the parsed JSON for potential structured queries
            filename = os.path.basename(file_path)
            self.loaded_dataframes[filename] = {"json_data": data}
            
            return json.dumps(data, indent=2)
        except Exception as e:
            logger.error(f"Error processing JSON {file_path}: {e}")
            return ""
    
    def _process_csv(self, file_path: str) -> List[Dict[str, Any]]:
        """Process CSV file with enhanced metadata."""
        logger.info(f"Processing CSV file: {file_path}")
        all_contents = []
        
        try:
            # Load CSV as DataFrame
            df = pd.read_csv(file_path)
            
            # Store the dataframe for potential querying
            filename = os.path.basename(file_path)
            self.loaded_dataframes[filename] = {"data": df}
            
            # Add the CSV content
            all_contents.append({
                "type": "table",
                "content": f"CSV file content:\n{df.to_string()}",
                "source": file_path
            })
            
            # Add DataFrame information for structured querying
            buffer = io.StringIO()
            df.info(buf=buffer)
            df_info = buffer.getvalue()
            
            all_contents.append({
                "type": "metadata",
                "content": f"CSV file structure:\nColumns: {', '.join(df.columns)}\n"
                           f"Shape: {df.shape[0]} rows, {df.shape[1]} columns\n"
                           f"Data types:\n{df_info}",
                "source": file_path
            })
            
            # Add summary statistics
            try:
                desc_stats = df.describe().to_string()
                all_contents.append({
                    "type": "statistics",
                    "content": f"CSV summary statistics:\n{desc_stats}",
                    "source": file_path
                })
            except Exception as e:
                logger.warning(f"Could not generate statistics for CSV: {e}")
                
        except Exception as e:
            logger.error(f"Error processing CSV {file_path}: {e}")
            
        return all_contents
    
    # def _crawl_links(self) -> List[Dict[str, Any]]:
    #     """Crawl extracted links to gather additional content."""
    #     logger.info(f"Crawling {len(self.extracted_links)} extracted links...")
    #     all_contents = []
        
    #     if not self.crawl_links or not self.extracted_links:
    #         return all_contents
        
    #     links_crawled = 0
    #     for link_info in self.extracted_links:
    #         # Extract URL from link info
    #         match = re.search(r'-> (.+)$', link_info)
    #         if not match:
    #             continue
                
    #         url = match.group(1).strip()
            
    #         # Skip non-HTTP links
    #         if not url.startswith(('http://', 'https://')):
    #             continue
                
    #         # Limit the number of links to crawl
    #         if links_crawled >= self.max_links_to_crawl:
    #             logger.info(f"Reached maximum links to crawl ({self.max_links_to_crawl})")
    #             break
                
    #         try:
    #             logger.info(f"Crawling link: {url}")
                
    #             # Fetch the web page with a timeout
    #             response = requests.get(url, timeout=10)
    #             response.raise_for_status()  # Raise an exception for bad responses
                
    #             # Parse HTML content
    #             soup = BeautifulSoup(response.text, 'html.parser')
                
    #             # Extract text content (remove script and style elements)
    #             for script in soup(["script", "style"]):
    #                 script.extract()
                
    #             # Get text and clean it up
    #             text = soup.get_text()
    #             lines = (line.strip() for line in text.splitlines())
    #             chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
    #             text = '\n'.join(chunk for chunk in chunks if chunk)
                
    #             if text:
    #                 all_contents.append({
    #                     "type": "web_content",
    #                     "content": f"Content from link: {url}\n\n{text}",
    #                     "source": url
    #                 })
                    
    #                 links_crawled += 1
                
    #             # Rate limiting to be polite
    #             time.sleep(1)
                
    #         except Exception as e:
    #             logger.warning(f"Failed to crawl link {url}: {e}")
        
    #     logger.info(f"Successfully crawled {links_crawled} links")
    #     return all_contents
    
    def _handle_structured_data_query(self, query: str) -> Optional[str]:
        """Process structured data queries against loaded dataframes."""
        if not self.loaded_dataframes:
            return None
            
        # Keywords that suggest a structured data query
        structured_keywords = [
            "average", "mean", "max", "min", "sum", "count", "how many", 
            "total", "percentage", "trend", "compare", "difference", 
            "higher than", "lower than", "calculate", "median",
            "top", "bottom", "largest", "smallest"
        ]
        
        # Check if this is likely a structured data query
        if not any(keyword in query.lower() for keyword in structured_keywords):
            return None
            
        # Generate responses for each dataframe
        responses = []
        
        for filename, dataframes in self.loaded_dataframes.items():
            for sheet_name, df in dataframes.items():
                if not isinstance(df, pd.DataFrame):
                    continue
                    
                response = f"Analysis of {filename}"
                if sheet_name != "data":
                    response += f" (sheet: {sheet_name})"
                response += ":\n"
                
                # Basic dataframe analysis based on query keywords
                if any(kw in query.lower() for kw in ["average", "mean"]):
                    try:
                        numeric_cols = df.select_dtypes(include=[np.number]).columns
                        if len(numeric_cols) > 0:
                            means = df[numeric_cols].mean()
                            response += f"Averages:\n{means.to_string()}\n\n"
                    except Exception as e:
                        logger.warning(f"Error calculating means: {e}")
                
                if any(kw in query.lower() for kw in ["max", "maximum", "highest", "largest", "top"]):
                    try:
                        numeric_cols = df.select_dtypes(include=[np.number]).columns
                        if len(numeric_cols) > 0:
                            maxes = df[numeric_cols].max()
                            response += f"Maximum values:\n{maxes.to_string()}\n\n"
                    except Exception as e:
                        logger.warning(f"Error calculating maximums: {e}")
                
                if any(kw in query.lower() for kw in ["min", "minimum", "lowest", "smallest", "bottom"]):
                    try:
                        numeric_cols = df.select_dtypes(include=[np.number]).columns
                        if len(numeric_cols) > 0:
                            mins = df[numeric_cols].min()
                            response += f"Minimum values:\n{mins.to_string()}\n\n"
                    except Exception as e:
                        logger.warning(f"Error calculating minimums: {e}")
                
                if any(kw in query.lower() for kw in ["sum", "total"]):
                    try:
                        numeric_cols = df.select_dtypes(include=[np.number]).columns
                        if len(numeric_cols) > 0:
                            sums = df[numeric_cols].sum()
                            response += f"Sum totals:\n{sums.to_string()}\n\n"
                    except Exception as e:
                        logger.warning(f"Error calculating sums: {e}")
                
                if any(kw in query.lower() for kw in ["count", "how many"]):
                    try:
                        response += f"Row count: {len(df)}\n"
                        response += f"Value counts:\n"
                        for col in df.columns[:3]:  # Limit to first 3 columns to avoid excessive output
                            try:
                                counts = df[col].value_counts().head(5)
                                response += f"\n{col}:\n{counts.to_string()}\n"
                            except:
                                pass
                    except Exception as e:
                        logger.warning(f"Error calculating counts: {e}")
                
                responses.append(response)
        
        if responses:
            return "\n\n".join(responses)
        return None
    
    def process_file(self, file_path: str) -> List[Document]:
        """Process a single file and return LlamaIndex Documents."""
        logger.info(f"Processing file: {file_path}")
        documents = []
        
        file_lower = file_path.lower()
        file_extension = os.path.splitext(file_lower)[1]
        
        try:
            if file_extension == '.pdf':
                contents = self._process_pdf_with_ocr(file_path)
                for item in contents:
                    doc = Document(
                        text=item["content"],
                        metadata={
                            "source": file_path,
                            "type": item["type"],
                            "page": item.get("page", None)
                        }
                    )
                    documents.append(doc)
            
            elif file_extension in ['.docx', '.doc']:
                contents = self._process_docx(file_path)
                for item in contents:
                    doc = Document(
                        text=item["content"],
                        metadata={
                            "source": file_path,
                            "type": item["type"]
                        }
                    )
                    documents.append(doc)
            
            elif file_extension in ['.jpg', '.jpeg', '.png', '.tiff', '.bmp', '.gif']:
                if self.use_ocr:
                    text = self._process_image_with_ocr(file_path)
                    if text:
                        doc = Document(
                            text=text,
                            metadata={
                                "source": file_path,
                                "type": "ocr_text"
                            }
                        )
                        documents.append(doc)
            
            elif file_extension in ['.xlsx', '.xls']:
                texts = self._process_excel(file_path)
                for text in texts:
                    doc = Document(
                        text=text,
                        metadata={
                            "source": file_path,
                            "type": "spreadsheet"
                        }
                    )
                    documents.append(doc)
            
            elif file_extension in ['.pptx', '.ppt']:
                contents = self._process_powerpoint(file_path)
                for item in contents:
                    doc = Document(
                        text=item["content"],
                        metadata={
                            "source": file_path,
                            "type": item["type"],
                            "slide": item.get("slide", None)
                        }
                    )
                    documents.append(doc)
            
            elif file_extension == '.json':
                text = self._process_json(file_path)
                if text:
                    doc = Document(
                        text=text,
                        metadata={
                            "source": file_path,
                            "type": "json"
                        }
                    )
                    documents.append(doc)
            
            elif file_extension == '.csv':
                contents = self._process_csv(file_path)
                for item in contents:
                    doc = Document(
                        text=item["content"],
                        metadata={
                            "source": file_path,
                            "type": item["type"]
                        }
                    )
                    documents.append(doc)
            
            else:
                # For text files and other formats, use SimpleDirectoryReader
                logger.info(f"Using SimpleDirectoryReader for file: {file_path}")
                reader = SimpleDirectoryReader(input_files=[file_path])
                documents.extend(reader.load_data())
        
        except Exception as e:
            logger.error(f"Error processing file {file_path}: {e}", exc_info=True)
        
        return documents
    
    # def process_directory(self, directory: str = None) -> List[Document]:
    #     """Process all files in a directory."""
    #     if directory is None:
    #         directory = self.data_dir
            
    #     logger.info(f"Processing directory: {directory}")
    #     all_documents = []
        
    #     for root, _, files in os.walk(directory):
    #         for file in files:
    #             file_path = os.path.join(root, file)
    #             documents = self.process_file(file_path)
    #             all_documents.extend(documents)
        
    #     # After processing all files, crawl extracted links if enabled
    #     if self.crawl_links and self.extracted_links:
    #         web_contents = self._crawl_links()
    #         for item in web_contents:
    #             doc = Document(
    #                 text=item["content"],
    #                 metadata={
    #                     "source": item["source"],
    #                     "type": item["type"]
    #                 }
    #             )
    #             all_documents.append(doc)
        
    #     return all_documents
    
    def create_index(self, documents: List[Document]) -> VectorStoreIndex:
        """Create a vector store index from documents."""
        logger.info(f"Creating index from {len(documents)} documents")
        
        # Parse documents into nodes
        nodes = self.node_parser.get_nodes_from_documents(documents)
        
        # Create vector store index
        index = VectorStoreIndex(nodes)
        
        return index
    
    def query_index(self, query: str, index: VectorStoreIndex) -> str:
        """Query the index with standard vector retrieval."""
        logger.info(f"Querying index with: {query}")
        
        # First check if this is a structured data query
        structured_response = self._handle_structured_data_query(query)
        if structured_response:
            return structured_response
        
        # Create custom retriever with configurable parameters
        retriever = index.as_retriever(
            similarity_top_k=self.similarity_top_k,
            similarity_cutoff=self.similarity_cutoff
        )
        
        # Custom prompt template for querying
        query_template = PromptTemplate(
            """You are a helpful document assistant that provides accurate information based on the document context.
            
            Context information is below:
            ---------------------
            {context_str}
            ---------------------
            
            Given this information, please answer the query. If the answer cannot be found in the document context, 
            please state "I don't have enough information to answer this question." Do not make up information.
            
            Query: {query_str}
            Answer: """
        )
        
        # Create query engine with custom parameters
        # query_engine = index.as_query_engine(
        #     retriever=retriever,
        #     text_qa_template=query_template,
        #     node_postprocessors=[]
        # )
        query_engine = RetrieverQueryEngine.from_args(
            retriever=retriever,
            text_qa_template=query_template,
            node_postprocessors=[]
        )
        # Execute query
        response = query_engine.query(query)
        
        return str(response)


class DocumentQAApp:
    """Streamlit app for document QA."""
    
    def __init__(self):
        # Default session state setup
        if "processor" not in st.session_state:
            st.session_state.processor = None
        if "index" not in st.session_state:
            st.session_state.index = None
        if "documents" not in st.session_state:
            st.session_state.documents = None
        if "messages" not in st.session_state:
            st.session_state.messages = []
        
    def run(self):
        """Run the Streamlit app."""
        

        st.title("Advanced Document QA System")
        
        # Sidebar for configuration
        with st.sidebar:
            st.header("Configuration")
            
            # data_dir = st.text_input("Data Directory", value="./data")
            
            # Model selection
            embedding_model = st.selectbox(
                "Embedding Model",
                ["BAAI/bge-small-en-v1.5", "BAAI/bge-base-en-v1.5", "sentence-transformers/all-mpnet-base-v2","models/embedding-001"],
                index=0
            )
            
            llm_model = st.selectbox(
                "LLM Model",
                ["StabilityAI/stablelm-tuned-alpha-3b", "meta-llama/Llama-2-7b-chat-hf","models/gemini-1.5-flash"],
                index=0
            )
            
            # Advanced parameters
            st.subheader("Advanced Parameters")
            
            chunk_size = st.slider("Chunk Size", min_value=256, max_value=4096, value=1024, step=128,
                                help="Size of text chunks for processing. Larger chunks provide more context but may reduce precision.")
            
            chunk_overlap = st.slider("Chunk Overlap", min_value=0, max_value=512, value=100, step=32,
                                    help="Overlap between chunks to maintain context continuity across chunks.")
            
            similarity_top_k = st.slider("Top K Results", min_value=1, max_value=10, value=5, step=1,
                                        help="Number of similar documents to retrieve per query.")
            
            similarity_cutoff = st.slider("Similarity Cutoff", min_value=0.0, max_value=1.0, value=0.7, step=0.05,
                                        help="Minimum similarity score required for retrieved documents.")
            
            use_ocr = st.checkbox("Use OCR", value=True,
                                help="Enable Optical Character Recognition for images and PDF files.")
            
            crawl_links = st.checkbox("Crawl Links", value=True,
                                    help="Crawl and index content from links found in documents.")
            
            max_links_to_crawl = st.slider("Max Links to Crawl", min_value=0, max_value=20, value=5, step=1,
                                        help="Maximum number of links to crawl from documents.")
            
            temperature = st.slider("LLM Temperature", min_value=0.0, max_value=1.0, value=0.2, step=0.1,
                                help="Controls randomness of responses. Lower values are more focused.")
            
            structured_data_threshold = st.slider("Structured Data Threshold", min_value=1, max_value=10, value=5, step=1,
                                                help="Threshold for determining if a query is about structured data.")
            

            # Button to initialize processor    
            # Initialize processor button
            if st.button("Initialize Processor"):
                with st.spinner("Initializing document processor..."):
                    processor = DocumentProcessor(
                        embedding_model=embedding_model,
                        llm_model=llm_model,
                        chunk_size=chunk_size,
                        chunk_overlap=chunk_overlap,
                        use_ocr=use_ocr,
                        similarity_top_k=similarity_top_k,
                        similarity_cutoff=similarity_cutoff,
                        crawl_links=crawl_links,
                        temperature=temperature,
                        max_links_to_crawl=max_links_to_crawl,
                        structured_data_threshold=structured_data_threshold
                    )
                    st.session_state.processor = processor
                st.success("Processor initialized!")

        
        # Main content area
        st.header("Document Processing")
        
        # File uploader
        uploaded_files = st.file_uploader("Upload documents", accept_multiple_files=True)
        
        if uploaded_files:
            if not st.session_state.processor:
                st.warning("Please initialize the processor from the sidebar before uploading and processing documents.")
            else:
                
                if st.button("Process Files"):
                    with st.spinner("Processing documents..."):
                        # Create temp directory for uploaded files
                        temp_dir = tempfile.mkdtemp()
                        try:
                            # Save uploaded files to temp directory
                            file_paths = []
                            for uploaded_file in uploaded_files:
                                file_path = os.path.join(temp_dir, uploaded_file.name)
                                with open(file_path, "wb") as f:
                                    f.write(uploaded_file.getbuffer())
                                file_paths.append(file_path)
                            
                            # Process each file
                            all_documents = []
                            for file_path in file_paths:
                                documents = st.session_state.processor.process_file(file_path)
                                all_documents.extend(documents)
                            
                            st.session_state.documents = all_documents
                            
                            # Create index
                            if all_documents:
                                st.session_state.index = st.session_state.processor.create_index(all_documents)
                                st.success(f"Processed {len(all_documents)} document chunks and created index!")
                            else:
                                st.error("No documents were successfully processed.")
                        finally:
                            # Clean up temp directory
                            shutil.rmtree(temp_dir)
            
        # Process directory
        # st.header("Process Directory")
        # if not st.session_state.processor:
        #     st.warning("Please initialize the processor from the sidebar before processing the data directory.")
        # else:
        #     if st.button("Process Data Directory"):
        #         with st.spinner(f"Processing directory {st.session_state.processor.data_dir}..."):
        #             st.session_state.documents = st.session_state.processor.process_directory()
                    
        #             if st.session_state.documents:
        #                 st.session_state.index = st.session_state.processor.create_index(st.session_state.documents)
        #                 st.success(f"Processed {len(st.session_state.documents)} document chunks from directory and created index!")
        #             else:
        #                 st.error("No documents were found or successfully processed in the directory.")
        
        # Query interface
        st.header("Ask Questions")
        
        for message in st.session_state.messages:
            with st.chat_message(message["role"]):
                st.markdown(message["content"])
        
        # Accept user input
        if query := st.chat_input("Ask your question about the documents..."):
            # Add user message to chat history
            st.session_state.messages.append({"role": "user", "content": query})
            
            # Display user message in chat message container
            with st.chat_message("user"):
                st.markdown(query)
            
            if not st.session_state.processor:
                response = "Please initialize the processor from the sidebar first."
            elif not st.session_state.index:
                response = "Please process documents first by uploading files."
            else:
                with st.spinner("Thinking..."):
                    response = st.session_state.processor.query_index(query, st.session_state.index)
            
            # Display assistant response in chat message container
            with st.chat_message("assistant"):
                st.markdown(response)
            
            # Add assistant response to chat history
            st.session_state.messages.append({"role": "assistant", "content": response})
        
        if st.button("Reset Chat"):
            st.session_state.messages = []
            st.rerun()            
            st.success("All states reset!")
        # query = st.text_input("Enter your question")
        
        # if st.button("Search"):
        #     if not query:
        #         st.warning("Please enter a question before searching.")
        #     elif not st.session_state.processor:
        #         st.warning("Please initialize the processor from the sidebar before asking questions.")
        #     elif not st.session_state.index:
        #         st.warning("Please process documents first by uploading files before asking questions.")
        #     else:
        #         with st.spinner("Searching..."):
        #             response = st.session_state.processor.query_index(query, st.session_state.index)

        #             st.subheader("Answer")
        #             st.write(response)

        #             # Show sources
        #             # st.subheader("Sources")
        #             # sources_shown = set()
        #             # for node in st.session_state.index.as_retriever().retrieve(query):
        #             #     source = node.metadata.get("source", "Unknown")
        #             #     if source not in sources_shown:
        #             #         node_type = node.metadata.get("type", "Unknown")
        #             #         page = node.metadata.get("page", "N/A")

        #             #         source_info = f"- {source}"
        #             #         if page and page != "N/A":
        #             #             source_info += f" (Page {page})"
        #             #         if node_type and node_type != "Unknown":
        #             #             source_info += f" [{node_type}]"

        #             #         st.write(source_info)
        #             #         sources_shown.add(source)
        # if st.button("Reset All"):
        #     st.session_state.processor = None
        #     st.session_state.index = None
        #     st.session_state.documents = None
        #     st.success("All states reset!")

if __name__ == "__main__":
    app = DocumentQAApp()
    app.run()