"""
Office document processing functionality (DOCX, PPTX, Excel).
"""
import docx
import pandas as pd
from pptx import Presentation
from typing import List, Dict, Any
from pathlib import Path
import json

from ..processors.ocr_processor import OCRProcessor
from ..processors.web_crawler import WebCrawler
from ..utils.logger import get_logger

logger = get_logger(__name__)

class OfficeProcessor:
    """Handles processing of Microsoft Office documents."""
    
    def __init__(self, use_ocr: bool = True, crawl_links: bool = True, max_links: int = 5):
        self.use_ocr = use_ocr
        self.crawl_links = crawl_links
        self.max_links = max_links
        self.ocr_processor = OCRProcessor() if use_ocr else None
        self.web_crawler = WebCrawler(max_links=max_links) if crawl_links else None
    
    def process_docx(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Process a DOCX file and extract content.
        
        Args:
            file_path: Path to the DOCX file
            
        Returns:
            List of content dictionaries
        """
        logger.info(f"Processing DOCX: {file_path}")
        content_items = []
        
        try:
            doc = docx.Document(file_path)
            
            # Extract paragraphs
            full_text = []
            for para in doc.paragraphs:
                if para.text.strip():
                    full_text.append(para.text)
            
            if full_text:
                content_items.append({
                    'type': 'text',
                    'content': '\n'.join(full_text),
                    'source': f"{Path(file_path).name} - Document Text"
                })
            
            # Extract tables
            table_content = self._extract_docx_tables(doc, file_path)
            content_items.extend(table_content)
            
            # Extract links if enabled
            if self.crawl_links:
                link_content = self._extract_docx_links(doc, file_path)
                content_items.extend(link_content)
            
            logger.info(f"Successfully processed DOCX: {len(content_items)} content items")
            
        except Exception as e:
            logger.error(f"Error processing DOCX {file_path}: {str(e)}")
        
        return content_items
    
    def process_pptx(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Process a PPTX file and extract content.
        
        Args:
            file_path: Path to the PPTX file
            
        Returns:
            List of content dictionaries
        """
        logger.info(f"Processing PPTX: {file_path}")
        content_items = []
        
        try:
            prs = Presentation(file_path)
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = self._process_slide(slide, slide_num, file_path)
                content_items.extend(slide_content)
            
            logger.info(f"Successfully processed PPTX: {len(content_items)} content items")
            
        except Exception as e:
            logger.error(f"Error processing PPTX {file_path}: {str(e)}")
        
        return content_items
    
    def process_excel(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Process an Excel file and extract content.
        
        Args:
            file_path: Path to the Excel file
            
        Returns:
            List of content dictionaries
        """
        logger.info(f"Processing Excel: {file_path}")
        content_items = []
        
        try:
            # Read all sheets
            excel_file = pd.ExcelFile(file_path)
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                
                if not df.empty:
                    # Convert DataFrame to readable text
                    sheet_text = self._dataframe_to_text(df, sheet_name)
                    
                    content_items.append({
                        'type': 'structured_data',
                        'content': sheet_text,
                        'source': f"{Path(file_path).name} - Sheet: {sheet_name}"
                    })
            
            logger.info(f"Successfully processed Excel: {len(content_items)} sheets")
            
        except Exception as e:
            logger.error(f"Error processing Excel {file_path}: {str(e)}")
        
        return content_items
    
    def process_csv(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Process a CSV file and extract content.
        
        Args:
            file_path: Path to the CSV file
            
        Returns:
            List of content dictionaries
        """
        logger.info(f"Processing CSV: {file_path}")
        content_items = []
        
        try:
            df = pd.read_csv(file_path)
            
            if not df.empty:
                csv_text = self._dataframe_to_text(df, "CSV Data")
                
                content_items.append({
                    'type': 'structured_data',
                    'content': csv_text,
                    'source': f"{Path(file_path).name} - CSV Data"
                })
            
            logger.info(f"Successfully processed CSV: {df.shape[0]} rows, {df.shape[1]} columns")
            
        except Exception as e:
            logger.error(f"Error processing CSV {file_path}: {str(e)}")
        
        return content_items
    
    def process_json(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Process a JSON file and extract content.
        
        Args:
            file_path: Path to the JSON file
            
        Returns:
            List of content dictionaries
        """
        logger.info(f"Processing JSON: {file_path}")
        content_items = []
        
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            
            # Convert JSON to readable text
            json_text = self._json_to_text(data)
            
            content_items.append({
                'type': 'structured_data',
                'content': json_text,
                'source': f"{Path(file_path).name} - JSON Data"
            })
            
            logger.info(f"Successfully processed JSON file")
            
        except Exception as e:
            logger.error(f"Error processing JSON {file_path}: {str(e)}")
        
        return content_items
    
    def _extract_docx_tables(self, doc, file_path: str) -> List[Dict[str, Any]]:
        """Extract tables from DOCX document."""
        content_items = []
        
        try:
            for table_num, table in enumerate(doc.tables, 1):
                table_data = []
                
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                
                if table_data:
                    # Convert to DataFrame for consistent formatting
                    df = pd.DataFrame(table_data[1:], columns=table_data[0] if table_data else None)
                    table_text = self._dataframe_to_text(df, f"Table {table_num}")
                    
                    content_items.append({
                        'type': 'table',
                        'content': table_text,
                        'source': f"{Path(file_path).name} - Table {table_num}"
                    })
                    
        except Exception as e:
            logger.error(f"Error extracting DOCX tables: {str(e)}")
        
        return content_items
    
    def _extract_docx_links(self, doc, file_path: str) -> List[Dict[str, Any]]:
        """Extract and crawl links from DOCX document."""
        content_items = []
        
        if not self.web_crawler:
            return content_items
        
        try:
            # Extract hyperlinks from document
            urls = []
            for rel in doc.part.rels.values():
                if "hyperlink" in rel.reltype:
                    urls.append(rel.target_ref)
            
            if urls:
                crawled_content = self.web_crawler.crawl_multiple_urls(urls[:self.max_links])
                
                for crawled in crawled_content:
                    content_items.append({
                        'type': 'web_content',
                        'content': f"Title: {crawled['title']}\n\nContent: {crawled['content']}",
                        'source': f"{Path(file_path).name} - Link: {crawled['url']}"
                    })
                    
        except Exception as e:
            logger.error(f"Error extracting DOCX links: {str(e)}")
        
        return content_items
    
    def _process_slide(self, slide, slide_num: int, file_path: str) -> List[Dict[str, Any]]:
        """Process a single PowerPoint slide."""
        content_items = []
        
        try:
            # Extract text from shapes
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            
            if slide_text:
                content_items.append({
                    'type': 'text',
                    'content': '\n'.join(slide_text),
                    'source': f"{Path(file_path).name} - Slide {slide_num}"
                })
            
            # Extract tables from slide
            for shape in slide.shapes:
                if shape.has_table:
                    table_text = self._extract_pptx_table(shape.table, slide_num)
                    if table_text:
                        content_items.append({
                            'type': 'table',
                            'content': table_text,
                            'source': f"{Path(file_path).name} - Slide {slide_num} Table"
                        })
                        
        except Exception as e:
            logger.error(f"Error processing slide {slide_num}: {str(e)}")
        
        return content_items
    
    def _extract_pptx_table(self, table, slide_num: int) -> str:
        """Extract table content from PowerPoint slide."""
        try:
            table_data = []
            
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                table_data.append(row_data)
            
            if table_data:
                df = pd.DataFrame(table_data[1:], columns=table_data[0] if table_data else None)
                return self._dataframe_to_text(df, f"Slide {slide_num} Table")
            
        except Exception as e:
            logger.error(f"Error extracting PPTX table: {str(e)}")
        
        return ""
    
    def _dataframe_to_text(self, df: pd.DataFrame, title: str) -> str:
        """Convert DataFrame to readable text format."""
        try:
            text_parts = [f"=== {title} ==="]
            
            # Add basic statistics
            text_parts.append(f"Shape: {df.shape[0]} rows, {df.shape[1]} columns")
            
            # Add column names
            text_parts.append(f"Columns: {', '.join(df.columns.astype(str))}")
            
            # Add sample data (first few rows)
            text_parts.append("\nSample Data:")
            text_parts.append(df.head().to_string(index=False))
            
            # Add summary statistics for numeric columns
            numeric_cols = df.select_dtypes(include=['number']).columns
            if len(numeric_cols) > 0:
                text_parts.append("\nNumeric Summary:")
                text_parts.append(df[numeric_cols].describe().to_string())
            
            return '\n'.join(text_parts)
            
        except Exception as e:
            logger.error(f"Error converting DataFrame to text: {str(e)}")
            return str(df)
    
    def _json_to_text(self, data: Any, max_depth: int = 3, current_depth: int = 0) -> str:
        """Convert JSON data to readable text format."""
        try:
            if current_depth > max_depth:
                return "[Max depth reached]"
            
            if isinstance(data, dict):
                text_parts = []
                for key, value in data.items():
                    if isinstance(value, (dict, list)):
                        nested_text = self._json_to_text(value, max_depth, current_depth + 1)
                        text_parts.append(f"{key}: {nested_text}")
                    else:
                        text_parts.append(f"{key}: {value}")
                return "{\n" + ",\n".join(text_parts) + "\n}"
            
            elif isinstance(data, list):
                if len(data) > 10:  # Limit list size for readability
                    sample_items = data[:5] + ["..."] + data[-2:]
                else:
                    sample_items = data
                
                text_parts = []
                for item in sample_items:
                    if isinstance(item, (dict, list)):
                        nested_text = self._json_to_text(item, max_depth, current_depth + 1)
                        text_parts.append(nested_text)
                    else:
                        text_parts.append(str(item))
                
                return "[" + ", ".join(text_parts) + "]"
            
            else:
                return str(data)
                
        except Exception as e:
            logger.error(f"Error converting JSON to text: {str(e)}")
            return str(data)
    
    def __del__(self):
        """Cleanup resources."""
        if self.web_crawler:
            self.web_crawler.close()
